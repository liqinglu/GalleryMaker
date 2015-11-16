#!/usr/bin/env Python
# coding=utf-8

import wx
import os
from pptx import Presentation
from pptx.util import Cm
from datetime import date


class PBDirFrame(wx.Frame):
    def __init__(self, app):
        wx.Frame.__init__(self, None, -1, u'文件选择框', size=(250,500))
        self.app = app
        font = wx.Font(12, wx.DEFAULT, wx.NORMAL, wx.NORMAL, False, 'Courier New') 
        self.SetFont(font)
        
        self.list = wx.ListBox(self, -1, (0,0), (200,600), '', wx.LB_SINGLE)
        self.list.Bind(wx.EVT_LISTBOX_DCLICK, self.OnDClick)
        
        curdir = os.getcwd()
        os.chdir(curdir)
        self.LoadDir(curdir)
        
        self.Bind(wx.EVT_CLOSE, self.OnClose)
    
        self.Show()
    
    def OnClose(self, event):
        self.Destroy()
        self.app.Close()
    
    def OnDClick(self, event):
        if self.list.GetSelection()==0:
            path = os.getcwd()
            pathinfo = os.path.split(path)
            dir = pathinfo[0]
        else:
            dir = self.list.GetStringSelection()
            # print os.path.abspath(dir)
        
        if os.path.isdir(dir):
            self.LoadDir(dir)
        elif os.path.splitext(dir)[-1] == '.jpg':
            self.app.ShowImage(dir)

    def LoadDir(self, dir):
        if not os.path.isdir(dir):
            return

        self.list.Clear()
        self.list.Append('...')
        dirs = []
        jpgs = []
        nnjpgs = []
        for _dir in os.listdir(dir):
            if os.path.isdir(dir+os.path.sep+_dir):
                dirs.append(_dir)
            else:
                info = os.path.splitext(_dir)
                # print info[0]
                if info[-1] == '.jpg':
                    if info[0].isdigit():
                        jpgs.append(string.atoi(info[0]))
                    else:
                        nnjpgs.append(_dir)
        jpgs.sort()
        for _jpgs in jpgs:
            self.list.Append(str(_jpgs)+'.jpg')
        for _nnjpgs in nnjpgs:
            self.list.Append(_nnjpgs)
        for _dirs in dirs:
            self.list.Append(_dirs)
        os.chdir(dir)

    def GetNextImage(self):
        index = self.list.GetSelection()
        i = index
        while i+1 < self.list.GetCount():
            i += 1
            if os.path.splitext(self.list.GetString(i))[-1] == '.jpg':
                break
        if i < self.list.GetCount():
            index = i
        self.list.SetSelection(index)
        return self.list.GetStringSelection()

    def GetPreImage(self):
        index = self.list.GetSelection()
        i = index
        while i-1 > 0:
            i -= 1
            if os.path.splitext(self.list.GetString(i))[-1] == '.jpg':
                break
        if i > 0:
            index = i
        
        self.list.SetSelection(index)
        return self.list.GetStringSelection()


class PBPicFrame(wx.Frame):
    max_width = 600
    max_height = 600

    def __init__(self, app):
        wx.Frame.__init__(self, None, -1, u'图形显示', size=(600,600))  # style=wx.SIMPLE_BORDER)

        self.bmoved = False
        
        self.app = app
        # staticbitmap
        self.bmp = wx.StaticBitmap(self, 0, wx.NullBitmap, (0,0), (600,600))

        self.Bind(wx.EVT_MOUSEWHEEL, self.OnChangeImage)
        self.bmp.Bind(wx.EVT_LEFT_DOWN, self.OnLeftDown)
        self.bmp.Bind(wx.EVT_LEFT_UP, self.OnLeftUp)
        self.bmp.Bind(wx.EVT_MOTION, self.OnMotion)
        self.Bind(wx.EVT_KEY_DOWN, self.OnKeyDown)
        
        self.ShowFullScreen(True, style=wx.FULLSCREEN_ALL)
        self.screenwidth, self.screenheight = wx.DisplaySize()
        # self.panel = wx.Panel(self)
        self.nextbutton = wx.Button(self, label='next', pos=(self.screenwidth/4,
                                                                   self.screenheight-50), size=(40, 40))
        self.prevbutton = wx.Button(self, label='prev', pos=(self.screenwidth/2,
                                                                   self.screenheight-50), size=(40, 40))
        self.addpptbutton = wx.Button(self, label='ToPPT', pos=(self.screenwidth*3/4,
                                                                      self.screenheight-50), size=(40, 40))
        self.Bind(wx.EVT_BUTTON, self.onnext, self.nextbutton)
        self.Bind(wx.EVT_BUTTON, self.onprev, self.prevbutton)
        self.Bind(wx.EVT_BUTTON, self.onaddppt, self.addpptbutton)

        self.Hide()

    def onnext(self, event):
        self.app.ShowNextImage()

    def onprev(self, event):
        self.app.ShowPreImage()

    def onaddppt(self, event):
        self.app.addtoppt()

    def ShowImage(self, path):
        if os.path.splitext(path)[-1] != '.jpg':
            return
        self.bmppath = path
        image = wx.Image(path, wx.BITMAP_TYPE_JPEG)
        bmp = image.ConvertToBitmap()
        size = self.GetSize(bmp)
        # print "actual size : %s, %s" % (size[0], size[1])
        bmp = image.Scale(size[0], size[1]).ConvertToBitmap()
        self.bmp.SetSize(size)
        self.bmp.SetBitmap(bmp)
        self.Show()

    def GetSize(self, bmp):
        width = bmp.GetWidth()
        height = bmp.GetHeight()
        # print "bmp size : %s, %s" % (width, height)
        if width > self.max_width:
            height = height*self.max_width/width
            width = self.max_width
        if height > self.max_height:
            width = width*self.max_height/height
            height = self.max_height
        size = width, height
        return size
       
    def OnChangeImage(self, event):
        rotation = event.GetWheelRotation()
        if rotation < 0:
            self.app.ShowNextImage()
        else:
            self.app.ShowPreImage()
    
    def OnLeftDown(self, event):
        self.pos = event.GetX(), event.GetY()
        self.bmoved = True

    def OnLeftUp(self, event):
        self.bmoved = False

    def OnMotion(self, event):
        if not self.bmoved:
            return
        pos = event.GetX(), event.GetY()
        dx = pos[0]-self.pos[0]
        dy = pos[1]-self.pos[1]
        pos = self.bmp.GetPosition()
        pos = pos[0]+dx, pos[1]+dy
        self.bmp.SetPosition(pos)

    def OnKeyDown(self, event):
        keycode = event.GetKeyCode()
        print keycode
        if keycode == 49:
            self.SizeUp()
        elif keycode == 50:
            self.SizeDown()
        event.Skip()

    def SizeUp(self):
        self.max_width += 50
        self.max_height += 75
        self.ShowImage(self.bmppath)

    def SizeDown(self):
        self.max_width -= 50
        self.max_height -= 75
        self.ShowImage(self.bmppath)


class PBApp(wx.App):

    def __init__(self, redirect=False):
        wx.App.__init__(self, redirect)
        self.pptlist = []
        self.prs = Presentation()
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "For Fun"
        subtitle.text = date.isoformat(date.today())
        self.destdir = os.getcwd()

    def OnInit(self):
        self.dirframe = PBDirFrame(self)
        self.picframe = PBPicFrame(self)
        self.Bind(wx.EVT_KEY_DOWN, self.OnKeyDown)
        return True

    def ShowImage(self, path):
        # print 'showing app img', path
        self.picframe.ShowImage(path)
        self.picframe.SetFocus()

    def ShowNextImage(self):
        path = self.dirframe.GetNextImage()
        self.ShowImage(path)

    def ShowPreImage(self):
        path = self.dirframe.GetPreImage()
        self.ShowImage(path)

    def addtoppt(self):
        path = self.dirframe.list.GetStringSelection()
        absolutepath = os.path.abspath(path)
        self.append(absolutepath)

    def append(self, path):
        # print path
        if path not in self.pptlist:
            self.pptlist.append(path)

            next_slide_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(next_slide_layout)
            left = Cm(0.1)
            top = Cm(4)
            width = Cm(25)
            slide.shapes.add_picture(path, left, top, width)

    def OnKeyDown(self, event):
        keycode = event.GetKeyCode()
        # print keycode
        if keycode == 27:
            if self.picframe.IsShown():
                self.picframe.Hide()
            else:
                self.picframe.Show()

    def Close(self):
        self.prs.save(self.destdir+"\\"+"test.pptx")
        self.picframe.Close()


def main():
    app = PBApp()
    app.MainLoop()

if __name__ == '__main__':
    main()
